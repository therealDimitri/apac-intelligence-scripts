#!/usr/bin/env python3
"""
Generate branded APAC Compass Survey PPTX using rich design elements from the
official Altera_Library_2026 slide library.

V3 — Uses pre-designed template slides with brand elements:
  Slide 32:  Title topo (3 placeholders)
  Slide 133: KPI doughnut chart cards (4 charts + big numbers + descriptions)
  Slide 53:  Branded data table (alternating purple/white rows)
  Slide 75:  Four ideas with icons + accent bars (pacing)
  Slide 78:  Three rounded cards with pill headers (keep for next year)
  Slide 112: Four-phase process with numbered ovals (AI narrative)
  Slide 80:  Quote light (verbatim feedback)
  Slide 174: APAC closer

Usage:
  export $(cat .env.local | grep -v '^#' | xargs)
  python3 scripts/generate-compass-pptx.py [--year 2026]

Requires: pip3 install python-pptx requests lxml
Output:  public/exports/APAC-Compass-{year}-Survey-Results.pptx
"""
import argparse
import glob
import json
import os
import re
import sys
from collections import Counter
from datetime import datetime
from pathlib import Path

import requests
from pptx import Presentation
from pptx.chart.data import ChartData
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

# ── Brand Colours ────────────────────────────────────────────────────
INDIGO       = RGBColor(0x38, 0x33, 0x92)
INDIGO_LIGHT = RGBColor(0xED, 0xEB, 0xF5)
NAVY         = RGBColor(0x15, 0x17, 0x44)
TEAL         = RGBColor(0x00, 0xBB, 0xBA)
BLUE         = RGBColor(0x00, 0x76, 0xA2)
VIOLET       = RGBColor(0x70, 0x7C, 0xF1)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
GREY_100     = RGBColor(0xF3, 0xF4, 0xF6)
GREY_400     = RGBColor(0x9C, 0xA3, 0xAF)
GREY_500     = RGBColor(0x6B, 0x72, 0x80)
GREY_700     = RGBColor(0x37, 0x41, 0x51)
EMERALD_700  = RGBColor(0x04, 0x78, 0x57)
AMBER        = RGBColor(0xF5, 0x9E, 0x0B)

FONT = 'Montserrat'

# ── Slide indices to KEEP from 174-slide template (0-based) ──────────
KEEP_SLIDES = {
    31:  'title',       # Slide 32:  Title topo
    132: 'kpi',         # Slide 133: 4 doughnut chart cards
    52:  'sessions',    # Slide 53:  Branded table
    74:  'pacing',      # Slide 75:  4 ideas with icons
    77:  'keep',        # Slide 78:  3 rounded cards
    111: 'narrative',   # Slide 112: 4-phase process
    79:  'verbatim',    # Slide 80:  Quote light
    173: 'closer',      # Slide 174: APAC closer
}

SLIDE_ORDER = [
    'title',
    'kpi',
    'sessions',
    'outcomes',         # programmatic
    'pacing',
    'keep',
    'narrative',
    'recommendations',  # programmatic (if AI data exists)
    'verbatim',
    'closer',
]


# ── Supabase Fetch ───────────────────────────────────────────────────

def fetch_data(year):
    url = os.environ.get('NEXT_PUBLIC_SUPABASE_URL')
    key = os.environ.get('SUPABASE_SERVICE_ROLE_KEY')
    if not url or not key:
        print("ERROR: Set NEXT_PUBLIC_SUPABASE_URL and SUPABASE_SERVICE_ROLE_KEY", file=sys.stderr)
        sys.exit(1)
    headers = {'apikey': key, 'Authorization': f'Bearer {key}'}
    r = requests.get(f'{url}/rest/v1/compass_surveys?year=eq.{year}&select=*', headers=headers)
    r.raise_for_status()
    surveys = r.json()
    if not surveys:
        print(f"ERROR: No survey found for year {year}", file=sys.stderr)
        sys.exit(1)
    survey = surveys[0]
    r2 = requests.get(f'{url}/rest/v1/compass_responses?survey_id=eq.{survey["id"]}&select=*&order=created_at', headers=headers)
    r2.raise_for_status()
    return {'survey': survey, 'responses': r2.json()}


# ── Shape Helpers ────────────────────────────────────────────────────

def delete_slide(prs, slide_idx):
    sldIdLst = prs.slides._sldIdLst
    sldId = sldIdLst[slide_idx]
    rId = sldId.get(qn('r:id'))
    prs.part.drop_rel(rId)
    sldIdLst.remove(sldId)


def is_placeholder(shape):
    try:
        shape.placeholder_format
        return True
    except ValueError:
        return False


def delete_shape(shape):
    sp = shape._element
    sp.getparent().remove(sp)


def rewrite_text(shape, lines, font_size=10, colour=NAVY, bold_first=True, alignment=None):
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    tf.clear()
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        if alignment:
            p.alignment = alignment
        run = p.add_run()
        run.text = line
        run.font.name = FONT
        run.font.size = Pt(font_size)
        run.font.color.rgb = colour
        run.font.bold = (bold_first and i == 0)


def set_placeholder(slide, idx, text, font_size=None, bold=None, colour=None):
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == idx:
            tf = ph.text_frame
            if tf.paragraphs and tf.paragraphs[0].runs:
                run = tf.paragraphs[0].runs[0]
                run.text = text
                if font_size: run.font.size = Pt(font_size)
                if bold is not None: run.font.bold = bold
                if colour: run.font.color.rgb = colour
            else:
                tf.clear()
                run = tf.paragraphs[0].add_run()
                run.text = text
                run.font.name = FONT
                if font_size: run.font.size = Pt(font_size)
                if bold is not None: run.font.bold = bold
                if colour: run.font.color.rgb = colour
            return


def style_cell(cell, text, font_size=9, bold=False, colour=NAVY, bg=None, alignment=PP_ALIGN.CENTER):
    cell.text = ''
    p = cell.text_frame.paragraphs[0]
    p.alignment = alignment
    run = p.add_run()
    run.text = str(text)
    run.font.size = Pt(font_size)
    run.font.name = FONT
    run.font.bold = bold
    run.font.color.rgb = colour
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    if bg:
        tcPr = cell._tc.get_or_add_tcPr()
        solidFill = tcPr.makeelement(qn('a:solidFill'), {})
        srgbClr = solidFill.makeelement(qn('a:srgbClr'), {'val': str(bg)})
        solidFill.append(srgbClr)
        tcPr.append(solidFill)


# ── Text Formatting ──────────────────────────────────────────────────

def clean_session(raw):
    return re.sub(r'^Rate your experience.*?:\s*', '', raw, flags=re.IGNORECASE).strip()[:45]

def clean_outcome(raw):
    return re.sub(r'^Strategic Outcomes?:\s*', '', raw, flags=re.IGNORECASE).strip()[:45]

def clean_venue(raw):
    return re.sub(r'^Venue/?s?:\s*', '', raw, flags=re.IGNORECASE).strip()[:40]

def clean_day(raw):
    m = re.match(r'Day\s*(\d)', raw, re.IGNORECASE)
    return f'Day {m.group(1)}' if m else raw.strip()[:20]

def pct(count, total):
    return f'{round(count / total * 100)}%' if total > 0 else '0%'

def nps_label(score):
    if score is None: return 'N/A'
    return f'+{score}' if score >= 0 else str(score)

def format_date_range(start, end):
    if not start: return ''
    months = ['January', 'February', 'March', 'April', 'May', 'June',
              'July', 'August', 'September', 'October', 'November', 'December']
    s = datetime.fromisoformat(start) if isinstance(start, str) else start
    if not end: return f'{s.day} {months[s.month - 1]} {s.year}'
    e = datetime.fromisoformat(end) if isinstance(end, str) else end
    if s.month == e.month: return f'{s.day}–{e.day} {months[s.month - 1]} {s.year}'
    return f'{s.day} {months[s.month - 1]} – {e.day} {months[e.month - 1]} {e.year}'

def collect_keys(responses, field):
    return list(dict.fromkeys(k for r in responses for k in (r.get(field) or {}).keys()))

def categorise_pacing(val):
    v = (val or '').lower()
    if 'perfect' in v: return 'Perfect'
    if 'too much' in v or 'too fast' in v or 'slightly too much' in v: return 'Too Much'
    if 'not enough' in v or 'too slow' in v: return 'Not Enough'
    if "didn't attend" in v or 'did not attend' in v: return "Didn't Attend"
    if v.strip(): return 'Other'
    return ''


# ══════════════════════════════════════════════════════════════════════
# SLIDE BUILDERS — Template Slide Population
# ══════════════════════════════════════════════════════════════════════

def fill_title(slide, survey, total):
    """Slide 32: Title topo — 3 placeholders."""
    location = survey.get('location') or ''
    dates = format_date_range(survey.get('date_start'), survey.get('date_end'))
    set_placeholder(slide, 0, f'APAC Compass {survey["year"]}', font_size=36, bold=False)
    set_placeholder(slide, 1, 'Survey Results & AI Intelligence', font_size=18, bold=False)
    set_placeholder(slide, 10, f'{location}  ·  {dates}  ·  {total} Responses', font_size=12)


def fill_kpi(slide, survey, responses):
    """Slide 133: 4 doughnut chart cards with big number overlays.
    Shapes: 4 charts (L→R), TextBox 24-27 (big nums at y<2.3"),
    TextBox 5-7,28 (details at y>=2.3")."""
    total = len(responses)
    nps = survey.get('nps_score') or 0
    mean_rec = float(survey.get('mean_recommend') or 0)
    mid_year = sum(1 for r in responses if r.get('mid_year_interest'))
    promoters = sum(1 for r in responses if (r.get('recommend_score') or 0) >= 9)
    passives = sum(1 for r in responses if 7 <= (r.get('recommend_score') or 0) <= 8)
    detractors = sum(1 for r in responses if (r.get('recommend_score') or 0) <= 6)

    set_placeholder(slide, 0, 'Key Metrics at a Glance', font_size=28, bold=False, colour=NAVY)

    # (chart_pct, big_text, detail_lines)
    kpis = [
        (max(0, min(100, round((nps + 100) / 2))),
         nps_label(nps),
         ['NPS Score', f'{promoters} Promoters · {passives} Passives · {detractors} Detractors']),
        (round(mean_rec * 10),
         f'{mean_rec:.1f}',
         ['Mean Recommend', 'Average score out of 10']),
        (round(mid_year / total * 100) if total else 0,
         pct(mid_year, total),
         ['Mid-Year Interest', f'{mid_year} of {total} want a mid-year event']),
        (100 if detractors == 0 else max(0, round((total - detractors) / total * 100)),
         str(detractors),
         ['Zero Detractors' if detractors == 0 else f'{detractors} Detractors',
          'All respondents scored 7+' if detractors == 0 else f'{detractors} of {total} scored 6 or below']),
    ]

    # Update doughnut chart data
    charts = sorted([s for s in slide.shapes if s.has_chart], key=lambda s: s.left)
    for i, cs in enumerate(charts[:4]):
        if i < len(kpis):
            try:
                cd = ChartData()
                cd.categories = ['Value', 'Remaining']
                cd.add_series('Data', (kpis[i][0], 100 - kpis[i][0]))
                cs.chart.replace_data(cd)
            except Exception as e:
                print(f'  Warning: chart {i+1} update failed: {e}')

    # Update text overlays (non-placeholder text boxes)
    tbs = [s for s in slide.shapes if s.has_text_frame and not is_placeholder(s)]
    big_nums = sorted([t for t in tbs if t.top < Inches(2.3)], key=lambda s: s.left)
    details = sorted([t for t in tbs if t.top >= Inches(2.3)], key=lambda s: s.left)

    for i, box in enumerate(big_nums[:4]):
        if i < len(kpis):
            rewrite_text(box, [kpis[i][1]], font_size=28, colour=NAVY, bold_first=True,
                         alignment=PP_ALIGN.CENTER)

    for i, box in enumerate(details[:4]):
        if i < len(kpis):
            rewrite_text(box, kpis[i][2], font_size=9, colour=GREY_700, bold_first=True,
                         alignment=PP_ALIGN.CENTER)


def fill_sessions(slide, responses):
    """Slide 53: Delete old 11×6 template table, add new branded table.
    Branded style: INDIGO header, alternating INDIGO_LIGHT/white rows."""
    total = len(responses)
    session_keys = collect_keys(responses, 'session_ratings')
    if not session_keys:
        return

    set_placeholder(slide, 0, 'Session Ratings', font_size=28, bold=False, colour=NAVY)

    # Set eyebrow placeholder (non-title placeholder)
    for ph in slide.placeholders:
        try:
            if ph.placeholder_format.idx != 0:
                tf = ph.text_frame
                if tf.paragraphs and tf.paragraphs[0].runs:
                    tf.paragraphs[0].runs[0].text = 'Survey Analysis'
                else:
                    tf.clear()
                    run = tf.paragraphs[0].add_run()
                    run.text = 'Survey Analysis'
                    run.font.name = FONT
                    run.font.size = Pt(10)
                break
        except ValueError:
            continue

    # Delete existing template table
    for shape in list(slide.shapes):
        if shape.shape_type == 19:  # MSO_SHAPE_TYPE.TABLE
            delete_shape(shape)

    # Build new branded table
    values = ['Very Relevant', 'Good', 'Not Relevant', "Didn't Attend"]
    rows = len(session_keys) + 1
    cols = len(values) + 1

    tbl = slide.shapes.add_table(
        rows, cols, Inches(0.3), Inches(1.4), Inches(9.4),
        Inches(min(0.30 * rows, 4.2))
    ).table

    tbl.columns[0].width = Inches(2.8)
    for ci in range(1, cols):
        tbl.columns[ci].width = Inches(6.6 // len(values))

    # Header
    style_cell(tbl.cell(0, 0), 'Session', 9, True, WHITE, INDIGO, PP_ALIGN.LEFT)
    for ci, v in enumerate(values):
        style_cell(tbl.cell(0, ci + 1), v, 8, True, WHITE, INDIGO)

    # Data rows
    for ri, key in enumerate(session_keys):
        counts = Counter(r.get('session_ratings', {}).get(key, '') for r in responses)
        bg = INDIGO_LIGHT if ri % 2 == 0 else None
        style_cell(tbl.cell(ri + 1, 0), clean_session(key), 8, False, NAVY, bg, PP_ALIGN.LEFT)
        for ci, v in enumerate(values):
            c = counts.get(v, 0)
            txt = f'{c} ({pct(c, total)})' if c > 0 else '–'
            style_cell(tbl.cell(ri + 1, ci + 1), txt, 8, c > 0, NAVY if c > 0 else GREY_400, bg)


def fill_pacing(slide, responses):
    """Slide 75: Four ideas with icons — Day 1, Day 2, Day 3, Venue.
    Shapes: Rectangle 3-6 (content cards, tall), Rectangle 7-10 (accent bars, thin),
    Picture 10-13 (icons, kept as-is)."""
    total = len(responses)
    set_placeholder(slide, 0, 'Day Pacing & Venue', font_size=28, bold=False, colour=NAVY)

    pacing_keys = collect_keys(responses, 'day_pacing')
    venue_keys = collect_keys(responses, 'venue_ratings')

    columns = []
    for key in pacing_keys[:3]:
        cats = [categorise_pacing(r.get('day_pacing', {}).get(key, '')) for r in responses]
        cat_counts = Counter(c for c in cats if c)
        day = clean_day(key)
        lines = [day, '']
        for cat in ['Perfect', 'Too Much', 'Not Enough', "Didn't Attend"]:
            c = cat_counts.get(cat, 0)
            if c > 0:
                lines.append(f'{cat}: {c} of {total}')
        columns.append(lines)

    # Venue summary as 4th column
    vlines = ['Venue Ratings', '']
    for key in venue_keys[:4]:
        counts = Counter(r.get('venue_ratings', {}).get(key, '') for r in responses)
        name = clean_venue(key)
        exc = counts.get('Excellent', 0)
        good = counts.get('Good', 0)
        vlines.append(f'{name}:')
        vlines.append(f'{exc + good}/{total} positive')
    columns.append(vlines)

    # Find tall content rectangles (not accent bars, not icons)
    rects = sorted(
        [s for s in slide.shapes
         if s.has_text_frame and not is_placeholder(s)
         and s.height > Inches(1.5) and s.width > Inches(1)],
        key=lambda s: s.left
    )

    for i, rect in enumerate(rects[:4]):
        if i < len(columns):
            rewrite_text(rect, columns[i], font_size=9, colour=NAVY, bold_first=True)


def fill_keep(slide, responses, year):
    """Slide 78: Three rounded cards with pill headers + takeaway boxes.
    Groups 2,5,8: each has Rounded Rectangle body + header pill.
    Rectangles 13-15: takeaway boxes at bottom."""
    keep_counts = Counter()
    for r in responses:
        for item in (r.get('keep_for_next_year') or []):
            if item.strip():
                keep_counts[item.strip()] += 1

    top3 = keep_counts.most_common(3)
    total = len(responses)

    set_placeholder(slide, 0, f'Keep for {year + 1}', font_size=28, bold=False, colour=NAVY)

    # Find and fill the 3 groups (sorted L→R)
    groups = sorted([s for s in slide.shapes if s.shape_type == 6], key=lambda s: s.left)

    for i, grp in enumerate(groups[:3]):
        if i >= len(top3):
            break
        item, count = top3[i]

        # Within each group: larger shape = body, smaller = header pill
        children = sorted(list(grp.shapes), key=lambda s: s.height, reverse=True)
        body = children[0] if children else None
        header = children[1] if len(children) > 1 else None

        if header and header.has_text_frame:
            rewrite_text(header, [item[:30]], font_size=10, colour=WHITE, bold_first=True,
                         alignment=PP_ALIGN.CENTER)

        if body and body.has_text_frame:
            body_lines = [
                f'{count} of {total} attendees ({pct(count, total)})',
                '',
                'Selected by majority' if count > total / 2 else f'Selected by {count} respondent{"s" if count != 1 else ""}',
            ]
            rewrite_text(body, body_lines, font_size=10, colour=NAVY, bold_first=True,
                         alignment=PP_ALIGN.CENTER)

    # Fill takeaway boxes (non-group rectangles below the cards)
    takeaways = sorted(
        [s for s in slide.shapes
         if s.has_text_frame and not is_placeholder(s)
         and s.shape_type != 6 and s.top > Inches(2.5) and s.width > Inches(1)],
        key=lambda s: s.left
    )

    extra = keep_counts.most_common(6)[3:]
    texts = [f'{item}: {cnt} of {total}' for item, cnt in extra]
    # Pad with summary stats if fewer than 3 extras
    while len(texts) < 3:
        summaries = [
            f'{total} responses analysed',
            f'{len(keep_counts)} unique items mentioned',
            f'Top: {top3[0][0]}' if top3 else 'No data',
        ]
        texts.append(summaries[len(texts)] if len(texts) < len(summaries) else '')

    for i, box in enumerate(takeaways[:3]):
        if i < len(texts):
            rewrite_text(box, [texts[i]], font_size=9, colour=NAVY, bold_first=False,
                         alignment=PP_ALIGN.CENTER)


def fill_narrative(slide, survey):
    """Slide 112: Four-phase process — numbered ovals + headings + body rects.
    Ovals 7-10: numbered circles (kept as 01-04).
    TextBox 2-5: heading text (rewrite with section titles).
    Rectangle 6, 11-13: body content (rewrite with narrative text)."""
    set_placeholder(slide, 0, 'AI Executive Narrative', font_size=28, bold=False, colour=NAVY)

    sections = [
        ('What Worked', 'Analysis pending'),
        ('The Surprise', 'Analysis pending'),
        ('The Risk', 'Analysis pending'),
        ('Quick Win', 'Analysis pending'),
    ]

    raw = survey.get('ai_narrative')
    if raw:
        try:
            data = json.loads(raw) if isinstance(raw, str) else raw
            sections = [
                ('What Worked', data.get('whatWorked', 'No data')),
                ('The Surprise', data.get('theSurprise', 'No data')),
                ('The Risk', data.get('theRisk', 'No data')),
                ('Quick Win', data.get('quickWin') or data.get('theQuickWin', 'No data')),
            ]
        except (json.JSONDecodeError, TypeError):
            pass

    non_ph = [s for s in slide.shapes if s.has_text_frame and not is_placeholder(s)]

    # Headings: TextBox shapes (short height), sorted L→R
    headings = sorted(
        [s for s in non_ph if 'TextBox' in s.name and s.height < Inches(1)],
        key=lambda s: s.left
    )
    # Bodies: Rectangle shapes (tall), sorted L→R
    bodies = sorted(
        [s for s in non_ph if 'Rectangle' in s.name and s.height > Inches(1)],
        key=lambda s: s.left
    )

    for i, h in enumerate(headings[:4]):
        if i < len(sections):
            rewrite_text(h, [sections[i][0]], font_size=11, colour=NAVY, bold_first=True,
                         alignment=PP_ALIGN.CENTER)

    for i, b in enumerate(bodies[:4]):
        if i < len(sections):
            text = sections[i][1][:200]
            rewrite_text(b, [text], font_size=9, colour=GREY_700, bold_first=False)


def fill_verbatim(slide, responses):
    """Slide 80: Quote light — single text placeholder for verbatim quotes."""
    quotes = []
    for r in sorted(responses, key=lambda r: r.get('recommend_score', 0), reverse=True):
        text = r.get('session_reasons') or r.get('change_one_thing') or r.get('other_comments') or ''
        if text.strip():
            name = r.get('respondent_name', 'Anonymous')
            score = r.get('recommend_score', '?')
            quotes.append(f'\u201c{text.strip()[:150]}\u201d\n\u2014 {name} (Score: {score})')

    target = None
    for ph in slide.placeholders:
        if ph.placeholder_format.idx == 10:
            target = ph
            break
    if not target:
        for ph in slide.placeholders:
            target = ph
            break
    if not target:
        return

    tf = target.text_frame
    tf.clear()
    for i, quote in enumerate(quotes[:4]):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.space_before = Pt(14) if i > 0 else Pt(0)
        run = p.add_run()
        run.text = quote
        run.font.name = FONT
        run.font.size = Pt(11)
        run.font.color.rgb = NAVY


# ══════════════════════════════════════════════════════════════════════
# SLIDE BUILDERS — Programmatic Tables (branded style)
# ══════════════════════════════════════════════════════════════════════

def add_outcomes_slide(prs, responses):
    """Programmatic branded table for strategic outcomes."""
    total = len(responses)
    keys = collect_keys(responses, 'outcome_ratings')
    if not keys:
        return None

    values = ['Strongly Agree', 'Agree', 'Neutral', 'Disagree', 'Strongly Disagree']
    layout = prs.slide_layouts[17]  # Title Only
    slide = prs.slides.add_slide(layout)
    set_placeholder(slide, 0, 'Strategic Outcomes', font_size=28, bold=False, colour=NAVY)

    rows = len(keys) + 1
    cols = len(values) + 1
    tbl = slide.shapes.add_table(
        rows, cols, Inches(0.3), Inches(1.2), Inches(9.4),
        Inches(min(0.45 * rows, 4.0))
    ).table

    tbl.columns[0].width = Inches(2.5)
    for ci in range(1, cols):
        tbl.columns[ci].width = Inches(6.9 // len(values))

    style_cell(tbl.cell(0, 0), 'Outcome', 9, True, WHITE, INDIGO, PP_ALIGN.LEFT)
    for ci, v in enumerate(values):
        style_cell(tbl.cell(0, ci + 1), v, 7, True, WHITE, INDIGO)

    for ri, key in enumerate(keys):
        counts = Counter(r.get('outcome_ratings', {}).get(key, '') for r in responses)
        bg = INDIGO_LIGHT if ri % 2 == 0 else None
        style_cell(tbl.cell(ri + 1, 0), clean_outcome(key), 8, False, NAVY, bg, PP_ALIGN.LEFT)
        for ci, v in enumerate(values):
            c = counts.get(v, 0)
            txt = f'{c} ({pct(c, total)})' if c > 0 else '–'
            style_cell(tbl.cell(ri + 1, ci + 1), txt, 8, False, NAVY if c > 0 else GREY_400, bg)

    return 'outcomes'


def add_recommendations_slide(prs, survey):
    """Programmatic branded table for AI recommendations."""
    raw = survey.get('ai_recommendations')
    if not raw:
        return None
    try:
        data = json.loads(raw) if isinstance(raw, str) else raw
    except (json.JSONDecodeError, TypeError):
        return None
    recs = data.get('recommendations', [])
    if not recs:
        return None

    layout = prs.slide_layouts[17]
    slide = prs.slides.add_slide(layout)
    set_placeholder(slide, 0, 'AI Recommendations', font_size=28, bold=False, colour=NAVY)

    rows = len(recs) + 1
    tbl = slide.shapes.add_table(
        rows, 4, Inches(0.3), Inches(1.2), Inches(9.4),
        Inches(min(0.50 * rows, 4.0))
    ).table

    tbl.columns[0].width = Inches(0.4)
    tbl.columns[1].width = Inches(5.5)
    tbl.columns[2].width = Inches(1.5)
    tbl.columns[3].width = Inches(2.0)

    for ci, h in enumerate(['#', 'Recommendation', 'Impact / Effort', 'Confidence']):
        style_cell(tbl.cell(0, ci), h, 9, True, WHITE, INDIGO,
                   PP_ALIGN.LEFT if ci == 1 else PP_ALIGN.CENTER)

    for ri, rec in enumerate(recs):
        bg = INDIGO_LIGHT if ri % 2 == 0 else None
        style_cell(tbl.cell(ri + 1, 0), str(ri + 1), 9, True, INDIGO, bg)
        style_cell(tbl.cell(ri + 1, 1), rec.get('recommendation', ''), 8, False, NAVY, bg, PP_ALIGN.LEFT)
        style_cell(tbl.cell(ri + 1, 2), f'{rec.get("impact", "")} / {rec.get("effort", "")}', 8, False, GREY_700, bg)
        conf = rec.get('confidence', '')
        cc = EMERALD_700 if 'high' in conf.lower() else (AMBER if 'med' in conf.lower() else GREY_500)
        style_cell(tbl.cell(ri + 1, 3), conf, 8, True, cc, bg)

    return 'recommendations'


# ══════════════════════════════════════════════════════════════════════
# POST-PROCESSING
# ══════════════════════════════════════════════════════════════════════

def remove_sections(prs):
    """Remove PowerPoint section markers from presentation XML."""
    for ext in prs._element.findall('.//' + qn('p:ext')):
        if '521415D9' in ext.get('uri', '').upper():
            parent = ext.getparent()
            parent.remove(ext)
            if len(parent) == 0:
                parent.getparent().remove(parent)
            break


# ══════════════════════════════════════════════════════════════════════
# MAIN
# ══════════════════════════════════════════════════════════════════════

if __name__ == '__main__':
    parser = argparse.ArgumentParser(description='Generate Compass Survey PPTX')
    parser.add_argument('--year', type=int, default=2026)
    parser.add_argument('--template', type=str, default=None)
    args = parser.parse_args()

    # Locate template
    template_path = args.template
    if not template_path:
        pattern = os.path.expanduser(
            '~/Library/CloudStorage/OneDrive-Altera*/**/Altera_Library_2026.pptx'
        )
        matches = glob.glob(pattern, recursive=True)
        if matches:
            template_path = matches[0]
        else:
            fb = '/tmp/altera-pptx-analysis/Altera_Library_2026.pptx'
            if os.path.exists(fb):
                template_path = fb

    if not template_path or not os.path.exists(template_path):
        print("ERROR: Cannot find Altera_Library_2026.pptx", file=sys.stderr)
        sys.exit(1)

    print(f'Template: {template_path}')
    print(f'Fetching data for {args.year}...')
    data = fetch_data(args.year)
    survey = data['survey']
    responses = data['responses']
    total = len(responses)
    print(f'  Survey: {survey["title"]}, {total} responses')

    # ── Open template ────────────────────────────────────────────────
    sz = os.path.getsize(template_path) / 1024 / 1024
    print(f'Opening template ({sz:.1f} MB)...')
    prs = Presentation(template_path)
    original = len(prs.slides)
    print(f'  {original} slides loaded')

    # ── Delete all slides except keepers ─────────────────────────────
    keep_set = set(KEEP_SLIDES.keys())
    for idx in range(original - 1, -1, -1):
        if idx not in keep_set:
            delete_slide(prs, idx)

    # Map remaining slides to roles (ascending original index order)
    slide_map = {}
    sorted_kept = sorted(KEEP_SLIDES.keys())
    for i, orig_idx in enumerate(sorted_kept):
        slide_map[KEEP_SLIDES[orig_idx]] = prs.slides[i]

    print(f'  Kept {len(prs.slides)} template slides, deleted {original - len(prs.slides)}')

    # ── Fill template slides ─────────────────────────────────────────
    print('Filling template slides...')
    fill_title(slide_map['title'], survey, total)
    fill_kpi(slide_map['kpi'], survey, responses)
    fill_sessions(slide_map['sessions'], responses)
    fill_pacing(slide_map['pacing'], responses)
    fill_keep(slide_map['keep'], responses, args.year)
    fill_narrative(slide_map['narrative'], survey)
    fill_verbatim(slide_map['verbatim'], responses)
    print('  All 7 template slides filled')

    # ── Add programmatic slides ──────────────────────────────────────
    print('Adding programmatic slides...')
    prog = []
    r = add_outcomes_slide(prs, responses)
    if r: prog.append(r)
    r = add_recommendations_slide(prs, survey)
    if r: prog.append(r)
    print(f'  Added {len(prog)} programmatic slides')

    # ── Reorder slides ───────────────────────────────────────────────
    print('Reordering slides...')
    current_roles = [KEEP_SLIDES[idx] for idx in sorted_kept] + prog

    desired = []
    for role in SLIDE_ORDER:
        if role in current_roles:
            desired.append(current_roles.index(role))

    sldIdLst = prs.slides._sldIdLst
    items = list(sldIdLst)
    for item in items:
        sldIdLst.remove(item)
    for idx in desired:
        sldIdLst.append(items[idx])

    print(f'  Final: {len(desired)} slides')

    # ── Remove PowerPoint sections ───────────────────────────────────
    remove_sections(prs)
    print('  Sections removed')

    # ── Save ─────────────────────────────────────────────────────────
    output_dir = Path(__file__).parent.parent / 'public' / 'exports'
    output_dir.mkdir(parents=True, exist_ok=True)
    output_path = output_dir / f'APAC-Compass-{args.year}-Survey-Results.pptx'
    prs.save(str(output_path))
    out_mb = output_path.stat().st_size / 1024 / 1024
    print(f'\nSaved: {output_path} ({out_mb:.1f} MB, {len(desired)} slides)')
