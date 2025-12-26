#!/usr/bin/env python3
"""
Generate PowerPoint Templates for APAC Client Success Intelligence Dashboard
Using Altera Digital Health branding
"""

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# Altera Brand Colors
ALTERA_PURPLE = RGBColor(124, 58, 237)  # #7C3AED - Primary purple
ALTERA_PURPLE_DARK = RGBColor(91, 33, 182)  # #5B21B6 - Dark purple
ALTERA_BLUE = RGBColor(59, 130, 246)  # #3B82F6 - Blue accent
ALTERA_GREEN = RGBColor(34, 197, 94)  # #22C55E - Green/success
ALTERA_ORANGE = RGBColor(249, 115, 22)  # #F97316 - Orange/warning
ALTERA_RED = RGBColor(239, 68, 68)  # #EF4444 - Red/critical
WHITE = RGBColor(255, 255, 255)
DARK_GRAY = RGBColor(55, 65, 81)  # #374151
LIGHT_GRAY = RGBColor(243, 244, 246)  # #F3F4F6

# Output directory
OUTPUT_DIR = os.path.join(os.path.dirname(os.path.dirname(__file__)), 'public', 'templates', 'presentations')


def create_title_slide(prs, title, subtitle="APAC Client Success"):
    """Create a title slide with Altera branding"""
    slide_layout = prs.slide_layouts[6]  # Blank layout
    slide = prs.slides.add_slide(slide_layout)

    # Purple header bar
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(2.5))
    header.fill.solid()
    header.fill.fore_color.rgb = ALTERA_PURPLE
    header.line.fill.background()

    # Title text
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.8), Inches(12), Inches(1))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = title
    title_para.font.size = Pt(44)
    title_para.font.bold = True
    title_para.font.color.rgb = WHITE

    # Subtitle
    subtitle_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.7), Inches(12), Inches(0.5))
    subtitle_frame = subtitle_box.text_frame
    subtitle_para = subtitle_frame.paragraphs[0]
    subtitle_para.text = subtitle
    subtitle_para.font.size = Pt(24)
    subtitle_para.font.color.rgb = RGBColor(233, 213, 255)  # Light purple

    # Footer with Altera branding
    footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.8), Inches(12), Inches(0.5))
    footer_frame = footer_box.text_frame
    footer_para = footer_frame.paragraphs[0]
    footer_para.text = "Altera Digital Health | APAC Client Success Intelligence Dashboard"
    footer_para.font.size = Pt(12)
    footer_para.font.color.rgb = DARK_GRAY

    # Placeholder for client/date
    client_box = slide.shapes.add_textbox(Inches(0.5), Inches(3.2), Inches(6), Inches(1.5))
    client_frame = client_box.text_frame
    p1 = client_frame.paragraphs[0]
    p1.text = "[Client Name]"
    p1.font.size = Pt(28)
    p1.font.bold = True
    p1.font.color.rgb = DARK_GRAY
    p2 = client_frame.add_paragraph()
    p2.text = "[Date]"
    p2.font.size = Pt(18)
    p2.font.color.rgb = RGBColor(107, 114, 128)  # Gray

    return slide


def create_section_slide(prs, section_title, color=ALTERA_PURPLE):
    """Create a section divider slide"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Full color background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = color
    bg.line.fill.background()

    # Section title
    title_box = slide.shapes.add_textbox(Inches(0.75), Inches(2.8), Inches(11.83), Inches(1.5))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = section_title
    title_para.font.size = Pt(48)
    title_para.font.bold = True
    title_para.font.color.rgb = WHITE
    title_para.alignment = PP_ALIGN.CENTER

    return slide


def create_content_slide(prs, title, bullet_points, notes=""):
    """Create a standard content slide with bullet points"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Purple header bar
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = ALTERA_PURPLE
    header.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(12), Inches(0.6))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = title
    title_para.font.size = Pt(28)
    title_para.font.bold = True
    title_para.font.color.rgb = WHITE

    # Content area
    content_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.6), Inches(12), Inches(5))
    text_frame = content_box.text_frame
    text_frame.word_wrap = True

    for i, point in enumerate(bullet_points):
        if i == 0:
            p = text_frame.paragraphs[0]
        else:
            p = text_frame.add_paragraph()
        p.text = f"• {point}"
        p.font.size = Pt(18)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(12)
        p.space_after = Pt(6)

    # Footer
    footer_box = slide.shapes.add_textbox(Inches(0.5), Inches(6.9), Inches(5), Inches(0.4))
    footer_frame = footer_box.text_frame
    footer_para = footer_frame.paragraphs[0]
    footer_para.text = "Altera Digital Health"
    footer_para.font.size = Pt(10)
    footer_para.font.color.rgb = RGBColor(156, 163, 175)

    # Page number placeholder
    page_box = slide.shapes.add_textbox(Inches(12), Inches(6.9), Inches(1), Inches(0.4))
    page_frame = page_box.text_frame
    page_para = page_frame.paragraphs[0]
    page_para.text = "[#]"
    page_para.font.size = Pt(10)
    page_para.font.color.rgb = RGBColor(156, 163, 175)
    page_para.alignment = PP_ALIGN.RIGHT

    return slide


def create_two_column_slide(prs, title, left_title, left_items, right_title, right_items):
    """Create a two-column layout slide"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Purple header bar
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = ALTERA_PURPLE
    header.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(12), Inches(0.6))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = title
    title_para.font.size = Pt(28)
    title_para.font.bold = True
    title_para.font.color.rgb = WHITE

    # Left column title
    left_title_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(5.8), Inches(0.5))
    ltf = left_title_box.text_frame
    ltp = ltf.paragraphs[0]
    ltp.text = left_title
    ltp.font.size = Pt(20)
    ltp.font.bold = True
    ltp.font.color.rgb = ALTERA_PURPLE

    # Left column content
    left_content = slide.shapes.add_textbox(Inches(0.5), Inches(2.1), Inches(5.8), Inches(4.5))
    left_frame = left_content.text_frame
    left_frame.word_wrap = True
    for i, item in enumerate(left_items):
        p = left_frame.paragraphs[0] if i == 0 else left_frame.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(8)

    # Right column title
    right_title_box = slide.shapes.add_textbox(Inches(6.8), Inches(1.5), Inches(5.8), Inches(0.5))
    rtf = right_title_box.text_frame
    rtp = rtf.paragraphs[0]
    rtp.text = right_title
    rtp.font.size = Pt(20)
    rtp.font.bold = True
    rtp.font.color.rgb = ALTERA_PURPLE

    # Right column content
    right_content = slide.shapes.add_textbox(Inches(6.8), Inches(2.1), Inches(5.8), Inches(4.5))
    right_frame = right_content.text_frame
    right_frame.word_wrap = True
    for i, item in enumerate(right_items):
        p = right_frame.paragraphs[0] if i == 0 else right_frame.add_paragraph()
        p.text = f"• {item}"
        p.font.size = Pt(16)
        p.font.color.rgb = DARK_GRAY
        p.space_before = Pt(8)

    return slide


def create_metrics_slide(prs, title, metrics):
    """Create a metrics/KPI slide with boxes"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Purple header bar
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = ALTERA_PURPLE
    header.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(12), Inches(0.6))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = title
    title_para.font.size = Pt(28)
    title_para.font.bold = True
    title_para.font.color.rgb = WHITE

    # Create metric boxes (up to 4)
    box_width = 2.8
    box_height = 1.8
    start_x = 0.75
    start_y = 1.8
    gap = 0.4

    colors = [ALTERA_PURPLE, ALTERA_BLUE, ALTERA_GREEN, ALTERA_ORANGE]

    for i, (metric_name, metric_value, metric_trend) in enumerate(metrics[:4]):
        x = start_x + (i * (box_width + gap))

        # Metric box
        box = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(start_y), Inches(box_width), Inches(box_height))
        box.fill.solid()
        box.fill.fore_color.rgb = colors[i % len(colors)]
        box.line.fill.background()

        # Metric value
        value_box = slide.shapes.add_textbox(Inches(x + 0.2), Inches(start_y + 0.3), Inches(box_width - 0.4), Inches(0.8))
        vf = value_box.text_frame
        vp = vf.paragraphs[0]
        vp.text = metric_value
        vp.font.size = Pt(36)
        vp.font.bold = True
        vp.font.color.rgb = WHITE
        vp.alignment = PP_ALIGN.CENTER

        # Metric name
        name_box = slide.shapes.add_textbox(Inches(x + 0.2), Inches(start_y + 1.1), Inches(box_width - 0.4), Inches(0.5))
        nf = name_box.text_frame
        np = nf.paragraphs[0]
        np.text = metric_name
        np.font.size = Pt(14)
        np.font.color.rgb = RGBColor(233, 213, 255)
        np.alignment = PP_ALIGN.CENTER

    # Trend/notes area
    if len(metrics) > 0:
        notes_box = slide.shapes.add_textbox(Inches(0.75), Inches(4.2), Inches(11.83), Inches(2.5))
        notes_frame = notes_box.text_frame
        notes_frame.word_wrap = True
        np = notes_frame.paragraphs[0]
        np.text = "Key Insights:"
        np.font.size = Pt(18)
        np.font.bold = True
        np.font.color.rgb = DARK_GRAY

        for metric_name, metric_value, metric_trend in metrics[:4]:
            p = notes_frame.add_paragraph()
            p.text = f"• {metric_name}: {metric_trend}"
            p.font.size = Pt(14)
            p.font.color.rgb = DARK_GRAY
            p.space_before = Pt(6)

    return slide


def create_table_slide(prs, title, headers, rows):
    """Create a slide with a table"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Purple header bar
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = ALTERA_PURPLE
    header.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(12), Inches(0.6))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = title
    title_para.font.size = Pt(28)
    title_para.font.bold = True
    title_para.font.color.rgb = WHITE

    # Create table
    num_cols = len(headers)
    num_rows = len(rows) + 1  # +1 for header
    col_width = 11.83 / num_cols

    table = slide.shapes.add_table(num_rows, num_cols, Inches(0.75), Inches(1.6), Inches(11.83), Inches(0.5 * num_rows)).table

    # Header row
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = ALTERA_PURPLE
        para = cell.text_frame.paragraphs[0]
        para.font.bold = True
        para.font.color.rgb = WHITE
        para.font.size = Pt(12)

    # Data rows
    for row_idx, row_data in enumerate(rows):
        for col_idx, cell_text in enumerate(row_data):
            cell = table.cell(row_idx + 1, col_idx)
            cell.text = str(cell_text)
            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(11)
            para.font.color.rgb = DARK_GRAY
            if row_idx % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_GRAY

    return slide


def create_action_items_slide(prs, title="Action Items & Next Steps"):
    """Create an action items slide with placeholders"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Purple header bar
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(1.2))
    header.fill.solid()
    header.fill.fore_color.rgb = ALTERA_PURPLE
    header.line.fill.background()

    # Title
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.35), Inches(12), Inches(0.6))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = title
    title_para.font.size = Pt(28)
    title_para.font.bold = True
    title_para.font.color.rgb = WHITE

    # Action items table
    headers = ["Action Item", "Owner", "Due Date", "Status"]
    table = slide.shapes.add_table(6, 4, Inches(0.75), Inches(1.6), Inches(11.83), Inches(2.5)).table

    # Set column widths
    table.columns[0].width = Inches(5)
    table.columns[1].width = Inches(2.5)
    table.columns[2].width = Inches(2)
    table.columns[3].width = Inches(2.33)

    # Header row
    for i, h in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = h
        cell.fill.solid()
        cell.fill.fore_color.rgb = ALTERA_PURPLE
        para = cell.text_frame.paragraphs[0]
        para.font.bold = True
        para.font.color.rgb = WHITE
        para.font.size = Pt(12)

    # Placeholder rows
    for row_idx in range(1, 6):
        for col_idx in range(4):
            cell = table.cell(row_idx, col_idx)
            cell.text = "[Enter details]" if col_idx == 0 else "[TBD]"
            para = cell.text_frame.paragraphs[0]
            para.font.size = Pt(11)
            para.font.color.rgb = RGBColor(156, 163, 175)
            if row_idx % 2 == 0:
                cell.fill.solid()
                cell.fill.fore_color.rgb = LIGHT_GRAY

    # Next meeting section
    next_box = slide.shapes.add_textbox(Inches(0.75), Inches(5), Inches(11.83), Inches(1.5))
    next_frame = next_box.text_frame
    np = next_frame.paragraphs[0]
    np.text = "Next Meeting"
    np.font.size = Pt(18)
    np.font.bold = True
    np.font.color.rgb = ALTERA_PURPLE

    p2 = next_frame.add_paragraph()
    p2.text = "Date: [Schedule next meeting]"
    p2.font.size = Pt(14)
    p2.font.color.rgb = DARK_GRAY
    p2.space_before = Pt(8)

    p3 = next_frame.add_paragraph()
    p3.text = "Focus: [Key topics for next session]"
    p3.font.size = Pt(14)
    p3.font.color.rgb = DARK_GRAY

    return slide


def create_closing_slide(prs, title="Thank You"):
    """Create a closing/thank you slide"""
    slide_layout = prs.slide_layouts[6]  # Blank
    slide = prs.slides.add_slide(slide_layout)

    # Full purple background
    bg = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(7.5))
    bg.fill.solid()
    bg.fill.fore_color.rgb = ALTERA_PURPLE
    bg.line.fill.background()

    # Thank you text
    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(2.5), Inches(12.33), Inches(1))
    title_frame = title_box.text_frame
    title_para = title_frame.paragraphs[0]
    title_para.text = title
    title_para.font.size = Pt(48)
    title_para.font.bold = True
    title_para.font.color.rgb = WHITE
    title_para.alignment = PP_ALIGN.CENTER

    # Contact info
    contact_box = slide.shapes.add_textbox(Inches(0.5), Inches(4), Inches(12.33), Inches(2))
    contact_frame = contact_box.text_frame

    cp1 = contact_frame.paragraphs[0]
    cp1.text = "[Your Name]"
    cp1.font.size = Pt(24)
    cp1.font.color.rgb = WHITE
    cp1.alignment = PP_ALIGN.CENTER

    cp2 = contact_frame.add_paragraph()
    cp2.text = "Client Success Engineer | Altera Digital Health"
    cp2.font.size = Pt(16)
    cp2.font.color.rgb = RGBColor(233, 213, 255)
    cp2.alignment = PP_ALIGN.CENTER

    cp3 = contact_frame.add_paragraph()
    cp3.text = "[email@alterahealth.com] | [Phone Number]"
    cp3.font.size = Pt(14)
    cp3.font.color.rgb = RGBColor(200, 180, 235)
    cp3.alignment = PP_ALIGN.CENTER
    cp3.space_before = Pt(12)

    return slide


# ============================================================================
# TEMPLATE GENERATORS
# ============================================================================

def generate_qbr_presentation():
    """Generate 6-slide QBR Presentation template"""
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    create_title_slide(prs, "Quarterly Business Review", "Q[X] [Year] Partnership Review")

    # Slide 2: Executive Summary
    create_content_slide(prs, "Executive Summary", [
        "Health Score: [XX]% ([+/-X]% from last quarter)",
        "NPS Score: [XX] ([Trend])",
        "Key Wins: [List 2-3 major achievements]",
        "Areas of Focus: [List 2-3 improvement areas]",
        "Overall Status: [Healthy/Watch/At Risk]"
    ])

    # Slide 3: Metrics Dashboard
    create_metrics_slide(prs, "Key Performance Metrics", [
        ("Health Score", "[XX]%", "[Describe trend and contributing factors]"),
        ("NPS", "[XX]", "[Describe sentiment changes]"),
        ("Meetings Held", "[XX]", "[Engagement frequency comparison]"),
        ("Actions Completed", "[XX]%", "[Task completion rate]")
    ])

    # Slide 4: Value Delivered
    create_two_column_slide(prs, "Value Delivered This Quarter",
        "Achievements", [
            "[Achievement 1 with metric]",
            "[Achievement 2 with metric]",
            "[Achievement 3 with metric]",
            "[Feature adoption milestone]"
        ],
        "ROI Highlights", [
            "Time saved: [XX hours/month]",
            "Efficiency gain: [XX]%",
            "Cost reduction: $[XX]",
            "Quality improvement: [Metric]"
        ]
    )

    # Slide 5: Success Plan
    create_two_column_slide(prs, "Next Quarter Success Plan",
        "Strategic Priorities", [
            "[Priority 1 - linked to business goal]",
            "[Priority 2 - linked to business goal]",
            "[Priority 3 - linked to business goal]"
        ],
        "Key Milestones", [
            "Month 1: [Specific milestone]",
            "Month 2: [Specific milestone]",
            "Month 3: [Specific milestone]"
        ]
    )

    # Slide 6: Action Items
    create_action_items_slide(prs)

    return prs


def generate_executive_business_review():
    """Generate Executive Business Review template"""
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # Slide 1: Title
    create_title_slide(prs, "Executive Business Review", "Strategic Partnership Assessment")

    # Slide 2: Strategic Alignment
    create_content_slide(prs, "Strategic Alignment", [
        "Business Objective 1: [Client goal] → Our Support: [How we help]",
        "Business Objective 2: [Client goal] → Our Support: [How we help]",
        "Business Objective 3: [Client goal] → Our Support: [How we help]",
        "Alignment Score: [Strong/Moderate/Needs Attention]"
    ])

    # Slide 3: Business Outcomes
    create_metrics_slide(prs, "Business Outcomes Achieved", [
        ("Revenue Impact", "$[XX]K", "[How we contributed to revenue]"),
        ("Efficiency Gain", "[XX]%", "[Operational improvements]"),
        ("Risk Reduction", "[XX]%", "[Compliance/security improvements]"),
        ("User Adoption", "[XX]%", "[Engagement metrics]")
    ])

    # Slide 4: Competitive Positioning
    create_two_column_slide(prs, "Competitive Positioning",
        "Our Strengths", [
            "[Unique value proposition 1]",
            "[Unique value proposition 2]",
            "[Industry expertise]",
            "[Support excellence]"
        ],
        "Market Comparison", [
            "[Benchmark vs industry]",
            "[Benchmark vs competitors]",
            "[Innovation leadership]",
            "[Partnership approach]"
        ]
    )

    # Slide 5: Joint Success Plan
    create_content_slide(prs, "Joint Success Plan: 12-Month Roadmap", [
        "Q1: [Major initiative and expected outcome]",
        "Q2: [Major initiative and expected outcome]",
        "Q3: [Major initiative and expected outcome]",
        "Q4: [Major initiative and expected outcome]",
        "Investment Required: [Resources, budget, time commitment]"
    ])

    # Slide 6: Executive Asks
    create_action_items_slide(prs, "Executive Action Items")

    return prs


def generate_checkin_meeting_agenda():
    """Generate Check-in Meeting Agenda template"""
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    # Single comprehensive slide
    slide_layout = prs.slide_layouts[6]
    slide = prs.slides.add_slide(slide_layout)

    # Header
    header = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(0), Inches(13.33), Inches(1))
    header.fill.solid()
    header.fill.fore_color.rgb = ALTERA_PURPLE
    header.line.fill.background()

    title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.25), Inches(8), Inches(0.5))
    tf = title_box.text_frame
    tp = tf.paragraphs[0]
    tp.text = "Check-in Meeting Agenda"
    tp.font.size = Pt(28)
    tp.font.bold = True
    tp.font.color.rgb = WHITE

    # Client/Date info
    info_box = slide.shapes.add_textbox(Inches(9), Inches(0.25), Inches(4), Inches(0.5))
    inf = info_box.text_frame
    ip = inf.paragraphs[0]
    ip.text = "[Client Name] | [Date]"
    ip.font.size = Pt(16)
    ip.font.color.rgb = RGBColor(233, 213, 255)
    ip.alignment = PP_ALIGN.RIGHT

    # Agenda items - Left column
    left_box = slide.shapes.add_textbox(Inches(0.5), Inches(1.3), Inches(6), Inches(5.8))
    left_frame = left_box.text_frame
    left_frame.word_wrap = True

    sections = [
        ("1. Recent Wins & Challenges (10 min)", [
            "Positive developments since last meeting",
            "Current challenges or blockers",
            "Support needed from Altera"
        ]),
        ("2. Product Usage Deep-Dive (10 min)", [
            "Usage metrics review",
            "Feature adoption status",
            "Training needs identification"
        ]),
        ("3. Open Support Items (5 min)", [
            "Outstanding tickets review",
            "Escalation status updates",
            "Upcoming maintenance/changes"
        ])
    ]

    for i, (section_title, items) in enumerate(sections):
        p = left_frame.paragraphs[0] if i == 0 else left_frame.add_paragraph()
        p.text = section_title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = ALTERA_PURPLE
        p.space_before = Pt(16) if i > 0 else Pt(0)

        for item in items:
            p = left_frame.add_paragraph()
            p.text = f"   • {item}"
            p.font.size = Pt(12)
            p.font.color.rgb = DARK_GRAY
            p.space_before = Pt(4)

    # Right column
    right_box = slide.shapes.add_textbox(Inches(6.8), Inches(1.3), Inches(6), Inches(5.8))
    right_frame = right_box.text_frame
    right_frame.word_wrap = True

    sections_right = [
        ("4. Action Item Review (5 min)", [
            "Previous action items status",
            "Blockers and dependencies",
            "Timeline adjustments needed"
        ]),
        ("5. Key Takeaways & Next Steps", [
            "Summary of decisions made",
            "New action items with owners",
            "Next meeting date confirmation"
        ])
    ]

    for i, (section_title, items) in enumerate(sections_right):
        p = right_frame.paragraphs[0] if i == 0 else right_frame.add_paragraph()
        p.text = section_title
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = ALTERA_PURPLE
        p.space_before = Pt(16) if i > 0 else Pt(0)

        for item in items:
            p = right_frame.add_paragraph()
            p.text = f"   • {item}"
            p.font.size = Pt(12)
            p.font.color.rgb = DARK_GRAY
            p.space_before = Pt(4)

    # Notes section
    notes_title = right_frame.add_paragraph()
    notes_title.text = "Meeting Notes:"
    notes_title.font.size = Pt(14)
    notes_title.font.bold = True
    notes_title.font.color.rgb = ALTERA_PURPLE
    notes_title.space_before = Pt(24)

    notes_box_shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(6.8), Inches(5), Inches(6), Inches(2))
    notes_box_shape.fill.solid()
    notes_box_shape.fill.fore_color.rgb = LIGHT_GRAY
    notes_box_shape.line.color.rgb = RGBColor(209, 213, 219)

    return prs


def generate_renewal_proposal():
    """Generate Renewal Proposal template"""
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    create_title_slide(prs, "Renewal Proposal", "Partnership Continuation & Growth")

    create_content_slide(prs, "Executive Summary", [
        "Partnership Duration: [X] years",
        "Current Contract Value: $[XXX,XXX]",
        "Proposed Renewal Term: [X] years",
        "Key Value Delivered: [Summary of impact]",
        "Recommendation: [Continue/Expand/Modify]"
    ])

    create_metrics_slide(prs, "12-Month Value Retrospective", [
        ("Total Value", "$[XXX]K", "[ROI calculation and methodology]"),
        ("Time Saved", "[XXX] hrs", "[Efficiency gains quantified]"),
        ("Quality Improvement", "[XX]%", "[Error reduction, compliance]"),
        ("User Satisfaction", "[XX]%", "[Survey results, NPS]")
    ])

    create_content_slide(prs, "ROI Analysis", [
        "Investment: $[XXX,XXX] annual subscription",
        "Return: $[XXX,XXX] in documented savings",
        "ROI: [XXX]%",
        "Payback Period: [X] months",
        "Intangible Benefits: [List qualitative improvements]"
    ])

    create_two_column_slide(prs, "Renewal Options",
        "Option A: Standard Renewal", [
            "Same scope and pricing",
            "[X]-year term",
            "Price: $[XXX,XXX]/year",
            "Includes: [List inclusions]"
        ],
        "Option B: Enhanced Package", [
            "Expanded scope with new modules",
            "[X]-year term with discount",
            "Price: $[XXX,XXX]/year",
            "Includes: [List additional features]"
        ]
    )

    create_action_items_slide(prs, "Renewal Next Steps")

    return prs


def generate_renewal_risk_assessment():
    """Generate Renewal Risk Assessment template"""
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    create_title_slide(prs, "Renewal Risk Assessment", "Strategic Analysis & Mitigation Plan")

    create_metrics_slide(prs, "Current Health Indicators", [
        ("Health Score", "[XX]%", "[Trend analysis]"),
        ("NPS", "[XX]", "[Sentiment analysis]"),
        ("Engagement", "[XX]%", "[Meeting/usage frequency]"),
        ("Support Load", "[XX]", "[Ticket volume trend]")
    ])

    create_two_column_slide(prs, "Stakeholder Mapping",
        "Key Decision Makers", [
            "[Name] - [Title] - [Sentiment]",
            "[Name] - [Title] - [Sentiment]",
            "[Name] - [Title] - [Sentiment]"
        ],
        "Influence Map", [
            "Executive Sponsor: [Name/Status]",
            "Day-to-Day Contact: [Name/Status]",
            "Technical Lead: [Name/Status]",
            "Procurement: [Name/Status]"
        ]
    )

    create_two_column_slide(prs, "Risk Analysis",
        "Competitive Threats", [
            "[Competitor 1]: [Threat level & response]",
            "[Competitor 2]: [Threat level & response]",
            "Internal alternatives: [Assessment]"
        ],
        "Risk Factors", [
            "Budget constraints: [High/Med/Low]",
            "Leadership change: [High/Med/Low]",
            "Product gaps: [High/Med/Low]",
            "Service issues: [High/Med/Low]"
        ]
    )

    create_content_slide(prs, "Mitigation Strategy", [
        "Immediate (This Week): [Action items]",
        "Short-Term (30 Days): [Action items]",
        "Medium-Term (60 Days): [Action items]",
        "Executive Escalation: [If needed, when and how]",
        "Success Criteria: [How we'll know risk is mitigated]"
    ])

    create_action_items_slide(prs, "Risk Mitigation Action Plan")

    return prs


def generate_multiyear_commitment():
    """Generate Multi-Year Commitment Pitch template"""
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    create_title_slide(prs, "Multi-Year Partnership", "Strategic Commitment Benefits")

    create_content_slide(prs, "Total Cost of Ownership Comparison", [
        "Year 1: $[XXX,XXX] (both options)",
        "Year 2: Annual = $[XXX,XXX] | Multi-Year = $[XXX,XXX]",
        "Year 3: Annual = $[XXX,XXX] | Multi-Year = $[XXX,XXX]",
        "3-Year Total: Annual = $[XXX,XXX] | Multi-Year = $[XXX,XXX]",
        "Savings with Multi-Year: $[XXX,XXX] ([XX]%)"
    ])

    create_two_column_slide(prs, "Multi-Year Discount Structure",
        "2-Year Commitment", [
            "[XX]% discount on Year 2",
            "Price lock guarantee",
            "Priority support included",
            "Annual price: $[XXX,XXX]"
        ],
        "3-Year Commitment", [
            "[XX]% discount on Years 2-3",
            "Price lock guarantee",
            "Premium support included",
            "Dedicated CSE assigned",
            "Annual price: $[XXX,XXX]"
        ]
    )

    create_content_slide(prs, "Commitment Benefits Package", [
        "✓ Price Protection: Locked rates for contract term",
        "✓ Priority Roadmap Input: Direct influence on development",
        "✓ Enhanced Support: [Specific support upgrades]",
        "✓ Training Credits: [XX] hours/year included",
        "✓ Executive Sponsor: Named executive relationship"
    ])

    create_content_slide(prs, "Payment Flexibility Options", [
        "Option 1: Annual payment (standard)",
        "Option 2: Quarterly payments (+[X]% admin fee)",
        "Option 3: Monthly payments (+[X]% admin fee)",
        "Option 4: Upfront payment (-[X]% discount)",
        "Custom arrangements available upon request"
    ])

    create_action_items_slide(prs, "Next Steps")

    return prs


def generate_upsell_discovery():
    """Generate Upsell Discovery Worksheet template"""
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    create_title_slide(prs, "Upsell Discovery", "Expansion Opportunity Analysis")

    create_content_slide(prs, "Current Product Usage Analysis", [
        "Active Modules: [List currently used modules]",
        "Usage Level: [High/Medium/Low] per module",
        "User Adoption: [XX]% of licensed users active",
        "Feature Utilisation: [XX]% of available features used",
        "Growth Pattern: [Describe usage trends]"
    ])

    create_two_column_slide(prs, "Opportunity Assessment",
        "Unutilised Features/Modules", [
            "[Feature 1] - [Potential value]",
            "[Feature 2] - [Potential value]",
            "[Feature 3] - [Potential value]",
            "[Module X] - [Potential value]"
        ],
        "Client Pain Points Addressed", [
            "[Pain point 1] → [Solution]",
            "[Pain point 2] → [Solution]",
            "[Pain point 3] → [Solution]"
        ]
    )

    create_content_slide(prs, "Business Case for Expansion", [
        "Current State: [Describe limitations]",
        "Proposed Addition: [Module/Feature name]",
        "Expected Benefits: [Quantified outcomes]",
        "Investment Required: $[XXX,XXX]",
        "Implementation Timeline: [X] weeks"
    ])

    create_metrics_slide(prs, "ROI Projection", [
        ("Investment", "$[XX]K", "[One-time + recurring costs]"),
        ("Annual Savings", "$[XX]K", "[Efficiency gains]"),
        ("ROI", "[XXX]%", "[Calculation methodology]"),
        ("Payback", "[X] mo", "[Break-even timeline]")
    ])

    create_action_items_slide(prs, "Discovery Next Steps")

    return prs


def generate_cross_sell_opportunity():
    """Generate Cross-Sell Opportunity Brief template"""
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    create_title_slide(prs, "Cross-Sell Opportunity", "Adjacent Solution Assessment")

    create_content_slide(prs, "Current Solution Footprint", [
        "Primary Product: [Product name and description]",
        "Contract Start: [Date]",
        "Annual Value: $[XXX,XXX]",
        "User Base: [XXX] users",
        "Satisfaction Level: [NPS/Health Score]"
    ])

    create_two_column_slide(prs, "Adjacent Products Alignment",
        "Recommended Product", [
            "[Product Name]",
            "Why it fits: [Alignment rationale]",
            "Integration: [How it connects]",
            "Price Point: $[XXX,XXX]/year"
        ],
        "Alternative Options", [
            "[Product 2]: [Brief description]",
            "[Product 3]: [Brief description]",
            "Bundle option: [Description]"
        ]
    )

    create_content_slide(prs, "Use Case Mapping", [
        "Current Workflow: [How they work today]",
        "Gap/Opportunity: [What's missing]",
        "Proposed Solution: [How new product helps]",
        "Expected Outcome: [Specific improvements]",
        "Success Metrics: [How we'll measure]"
    ])

    create_content_slide(prs, "Bundled Pricing Options", [
        "Standalone Price: $[XXX,XXX]/year",
        "Bundle Discount: [XX]% when combined with current products",
        "Bundled Price: $[XXX,XXX]/year",
        "Implementation: [Included/Additional cost]",
        "Training: [Included/Additional cost]"
    ])

    create_action_items_slide(prs, "Cross-Sell Next Steps")

    return prs


def generate_expansion_business_case():
    """Generate Expansion Business Case template"""
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    create_title_slide(prs, "Expansion Business Case", "Investment Justification")

    create_content_slide(prs, "Current State Analysis", [
        "Existing Investment: $[XXX,XXX]/year",
        "Current Capabilities: [List what's in place]",
        "Identified Gaps: [What's missing]",
        "Business Impact of Gaps: [Quantified if possible]",
        "Strategic Alignment: [How expansion supports goals]"
    ])

    create_content_slide(prs, "Quantified Benefits Breakdown", [
        "Direct Cost Savings: $[XXX,XXX]/year",
        "Productivity Gains: [XX] hours/week × [XX] staff = $[XXX,XXX]",
        "Revenue Enablement: $[XXX,XXX] potential",
        "Risk Reduction: [Quantified or qualitative]",
        "Total Annual Value: $[XXX,XXX]"
    ])

    create_content_slide(prs, "Investment Required", [
        "Software Licensing: $[XXX,XXX]/year",
        "Implementation Services: $[XXX,XXX] (one-time)",
        "Training: $[XXX,XXX] (one-time)",
        "Internal Resources: [XX] FTE for [XX] weeks",
        "Total Year 1 Investment: $[XXX,XXX]"
    ])

    create_metrics_slide(prs, "Payback Period Calculation", [
        ("Year 1", "$[XX]K", "[Net benefit after costs]"),
        ("Year 2", "$[XX]K", "[Ongoing annual benefit]"),
        ("Year 3", "$[XX]K", "[Ongoing annual benefit]"),
        ("Payback", "[X] mo", "[When investment is recovered]")
    ])

    create_content_slide(prs, "Risk Assessment & Mitigation", [
        "Implementation Risk: [Level] - Mitigation: [Plan]",
        "Adoption Risk: [Level] - Mitigation: [Plan]",
        "Budget Risk: [Level] - Mitigation: [Plan]",
        "Timeline Risk: [Level] - Mitigation: [Plan]"
    ])

    create_action_items_slide(prs, "Approval & Next Steps")

    return prs


def generate_at_risk_recovery_plan():
    """Generate At-Risk Client Recovery Plan template"""
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    create_title_slide(prs, "Client Recovery Plan", "Relationship Restoration Strategy")

    create_content_slide(prs, "Current Situation Assessment", [
        "Health Score: [XX]% (down from [XX]%)",
        "NPS: [XX] ([Detractor/Passive])",
        "Primary Issues: [List top 3 concerns]",
        "Relationship Status: [Critical/At Risk/Watch]",
        "Renewal Date: [Date] ([XX] days away)"
    ])

    create_content_slide(prs, "Root Cause Analysis", [
        "Issue 1: [Description] - Root Cause: [Finding]",
        "Issue 2: [Description] - Root Cause: [Finding]",
        "Issue 3: [Description] - Root Cause: [Finding]",
        "Contributing Factors: [Internal/External factors]",
        "Pattern Analysis: [Is this systemic or isolated?]"
    ])

    create_two_column_slide(prs, "Recovery Action Plan",
        "Quick Wins (30 Days)", [
            "[Action 1] - Owner: [Name]",
            "[Action 2] - Owner: [Name]",
            "[Action 3] - Owner: [Name]",
            "Success Metric: [How we'll know]"
        ],
        "Medium-Term (60-90 Days)", [
            "[Action 1] - Owner: [Name]",
            "[Action 2] - Owner: [Name]",
            "[Action 3] - Owner: [Name]",
            "Success Metric: [How we'll know]"
        ]
    )

    create_content_slide(prs, "Executive Escalation Strategy", [
        "Internal Escalation: [Who needs to be involved]",
        "Client Executive Engagement: [Who to contact]",
        "Escalation Trigger: [When to escalate further]",
        "Executive Sponsor Ask: [What we need from leadership]",
        "Communication Plan: [How we'll keep everyone informed]"
    ])

    create_action_items_slide(prs, "Recovery Milestones & Owners")

    return prs


def generate_attrition_prevention_playbook():
    """Generate Attrition Prevention Playbook template"""
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    create_title_slide(prs, "Attrition Prevention Playbook", "Retention Strategy Guide")

    create_content_slide(prs, "Early Warning Indicators", [
        "🔴 Critical: Health score drops >20 points in 30 days",
        "🔴 Critical: Executive sponsor leaves organization",
        "🟡 Warning: Support tickets increase >50%",
        "🟡 Warning: Meeting attendance declines",
        "🟢 Watch: Minor NPS decline (1-2 points)"
    ])

    create_two_column_slide(prs, "Intervention Tactics by Scenario",
        "Budget/Cost Concerns", [
            "Present ROI analysis",
            "Offer flexible payment terms",
            "Propose scope reduction vs cancellation",
            "Escalate to leadership with value proof"
        ],
        "Product/Service Issues", [
            "Acknowledge and apologise",
            "Create action plan with timeline",
            "Assign dedicated support contact",
            "Provide service credits if appropriate"
        ]
    )

    create_two_column_slide(prs, "Intervention Tactics (Continued)",
        "Competitive Pressure", [
            "Conduct competitive analysis",
            "Highlight switching costs/risks",
            "Present roadmap alignment",
            "Offer loyalty incentives"
        ],
        "Organizational Change", [
            "Map new stakeholders quickly",
            "Schedule executive introduction",
            "Re-establish value narrative",
            "Offer new team training"
        ]
    )

    create_content_slide(prs, "Win-Back Offer Structure", [
        "Tier 1 (High Value): [Discount %] + [Added services]",
        "Tier 2 (Medium Value): [Discount %] + [Added services]",
        "Tier 3 (Standard): [Discount %]",
        "Approval Required: [Who approves which tier]",
        "Time Limit: Offers valid for [XX] days"
    ])

    create_content_slide(prs, "Escalation Protocol", [
        "Day 1-3: CSE direct outreach and assessment",
        "Day 4-7: Manager involvement and action plan",
        "Day 8-14: Director escalation if no progress",
        "Day 15+: VP/Executive engagement for critical accounts",
        "Documentation: All interactions logged in [System]"
    ])

    create_action_items_slide(prs, "Playbook Action Items")

    return prs


def generate_discovery_call_template():
    """Generate Discovery Call Template"""
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    create_title_slide(prs, "Discovery Call", "Understanding Your Needs")

    create_content_slide(prs, "Business Objectives", [
        "What are your top 3 priorities for [Year/Quarter]?",
        "[Objective 1]: [Notes]",
        "[Objective 2]: [Notes]",
        "[Objective 3]: [Notes]",
        "How does [our solution area] fit into these priorities?"
    ])

    create_content_slide(prs, "Current Challenges", [
        "What's working well today?",
        "What are your biggest pain points?",
        "[Challenge 1]: [Impact and urgency]",
        "[Challenge 2]: [Impact and urgency]",
        "What happens if these aren't addressed?"
    ])

    create_content_slide(prs, "Decision Criteria & Stakeholders", [
        "What criteria will you use to evaluate solutions?",
        "Who else is involved in this decision?",
        "Decision Maker: [Name/Title]",
        "Influencers: [Names/Titles]",
        "What's their evaluation process?"
    ])

    create_content_slide(prs, "Budget & Timeline", [
        "Is there budget allocated for this initiative?",
        "Budget Range: $[XXX] - $[XXX]",
        "What's your ideal implementation timeline?",
        "Are there any deadlines driving this project?",
        "What would delay or accelerate the decision?"
    ])

    create_action_items_slide(prs, "Discovery Follow-Up Items")

    return prs


def generate_demo_customization():
    """Generate Demo Customization Guide template"""
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    create_title_slide(prs, "Demo Customisation Guide", "Tailored Presentation Plan")

    create_content_slide(prs, "Audience Profile", [
        "Company: [Name]",
        "Industry: [Healthcare/Government/etc.]",
        "Size: [XX] employees / $[XX]M revenue",
        "Attendees: [List names and titles]",
        "Technical Level: [High/Medium/Low]"
    ])

    create_content_slide(prs, "Key Pain Points to Address", [
        "Pain Point 1: [Description] → Demo Feature: [Feature name]",
        "Pain Point 2: [Description] → Demo Feature: [Feature name]",
        "Pain Point 3: [Description] → Demo Feature: [Feature name]",
        "Must-Show Feature: [Based on discovery]",
        "Avoid Showing: [Features that don't apply]"
    ])

    create_content_slide(prs, "Features to Highlight", [
        "✓ [Feature 1]: [Why it matters to this client]",
        "✓ [Feature 2]: [Why it matters to this client]",
        "✓ [Feature 3]: [Why it matters to this client]",
        "✓ Integration with: [Their existing systems]",
        "✗ Skip: [Features not relevant]"
    ])

    create_content_slide(prs, "ROI Talking Points", [
        "Similar Client Result: [Company] achieved [XX]% improvement",
        "Time Savings: Estimated [XX] hours/week",
        "Cost Reduction: Potential $[XXX,XXX] annual savings",
        "Risk Mitigation: [Compliance/security benefits]",
        "Competitive Advantage: [How it differentiates them]"
    ])

    create_action_items_slide(prs, "Demo Preparation Checklist")

    return prs


def generate_proposal_template():
    """Generate Proposal Template"""
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    create_title_slide(prs, "Solution Proposal", "Prepared for [Client Name]")

    create_content_slide(prs, "Executive Summary", [
        "Opportunity: [One sentence summary]",
        "Proposed Solution: [High-level description]",
        "Investment: $[XXX,XXX] [per year/one-time]",
        "Expected ROI: [XXX]% over [X] years",
        "Implementation: [X] weeks"
    ])

    create_content_slide(prs, "Proposed Solution", [
        "Core Components: [List main elements]",
        "[Component 1]: [Description and value]",
        "[Component 2]: [Description and value]",
        "[Component 3]: [Description and value]",
        "Optional Add-ons: [List available options]"
    ])

    create_content_slide(prs, "Implementation Plan", [
        "Phase 1 (Weeks 1-2): [Discovery and planning]",
        "Phase 2 (Weeks 3-4): [Configuration and setup]",
        "Phase 3 (Weeks 5-6): [Testing and training]",
        "Phase 4 (Week 7+): [Go-live and support]",
        "Your Team's Time: [XX] hours total"
    ])

    create_content_slide(prs, "Pricing & Terms", [
        "Software License: $[XXX,XXX]/year",
        "Implementation: $[XXX,XXX] (one-time)",
        "Training: $[XXX,XXX] (one-time)",
        "Total Year 1: $[XXX,XXX]",
        "Terms: [Net 30 / Annual prepaid / etc.]"
    ])

    create_action_items_slide(prs, "Proposal Next Steps")

    create_closing_slide(prs, "Thank You")

    return prs


def generate_negotiation_prep():
    """Generate Negotiation Prep Template"""
    prs = Presentation()
    prs.slide_width = Inches(13.33)
    prs.slide_height = Inches(7.5)

    create_title_slide(prs, "Negotiation Preparation", "Deal Strategy & Boundaries")

    create_content_slide(prs, "Walk-Away Price & Boundaries", [
        "List Price: $[XXX,XXX]",
        "Target Price: $[XXX,XXX] ([XX]% discount)",
        "Floor Price: $[XXX,XXX] ([XX]% max discount)",
        "Walk-Away Point: $[XXX,XXX]",
        "Approval Required Below: $[XXX,XXX]"
    ])

    create_two_column_slide(prs, "Concession Strategy",
        "What We Can Offer", [
            "Payment terms flexibility",
            "Extended implementation support",
            "Additional training hours",
            "Multi-year price lock",
            "Professional services discount"
        ],
        "What We Need in Return", [
            "Multi-year commitment",
            "Case study participation",
            "Reference availability",
            "Faster decision timeline",
            "Expanded scope"
        ]
    )

    create_content_slide(prs, "Value Reinforcement Points", [
        "Unique Differentiator 1: [What competitors can't match]",
        "Unique Differentiator 2: [What competitors can't match]",
        "Risk of Alternative: [What they lose with competitor]",
        "Switching Cost: [True cost of not choosing us]",
        "Success Story: [Relevant customer proof point]"
    ])

    create_content_slide(prs, "Close Plan", [
        "Decision Maker: [Name and their priorities]",
        "Timeline: Decision needed by [Date]",
        "Competitor Status: [Who else is being considered]",
        "Key Objection: [Anticipated pushback]",
        "Response: [How we'll address it]"
    ])

    create_action_items_slide(prs, "Pre-Negotiation Checklist")

    return prs


# ============================================================================
# MAIN EXECUTION
# ============================================================================

def main():
    """Generate all templates"""
    os.makedirs(OUTPUT_DIR, exist_ok=True)

    templates = [
        # Customer Meeting Templates
        ("QBR-Presentation-Template.pptx", generate_qbr_presentation, "QBR Presentation"),
        ("Executive-Business-Review-Template.pptx", generate_executive_business_review, "Executive Business Review"),
        ("Check-in-Meeting-Agenda-Template.pptx", generate_checkin_meeting_agenda, "Check-in Meeting Agenda"),

        # Renewal Templates
        ("Renewal-Proposal-Template.pptx", generate_renewal_proposal, "Renewal Proposal"),
        ("Renewal-Risk-Assessment-Template.pptx", generate_renewal_risk_assessment, "Renewal Risk Assessment"),
        ("Multi-Year-Commitment-Template.pptx", generate_multiyear_commitment, "Multi-Year Commitment"),

        # Expansion Templates
        ("Upsell-Discovery-Template.pptx", generate_upsell_discovery, "Upsell Discovery"),
        ("Cross-Sell-Opportunity-Template.pptx", generate_cross_sell_opportunity, "Cross-Sell Opportunity"),
        ("Expansion-Business-Case-Template.pptx", generate_expansion_business_case, "Expansion Business Case"),

        # Risk Mitigation Templates
        ("At-Risk-Recovery-Plan-Template.pptx", generate_at_risk_recovery_plan, "At-Risk Recovery Plan"),
        ("Attrition-Prevention-Playbook-Template.pptx", generate_attrition_prevention_playbook, "Attrition Prevention Playbook"),

        # Sales Process Templates
        ("Discovery-Call-Template.pptx", generate_discovery_call_template, "Discovery Call"),
        ("Demo-Customization-Template.pptx", generate_demo_customization, "Demo Customization"),
        ("Proposal-Template.pptx", generate_proposal_template, "Proposal"),
        ("Negotiation-Prep-Template.pptx", generate_negotiation_prep, "Negotiation Prep"),
    ]

    print("=" * 60)
    print("Generating Altera PowerPoint Templates")
    print("=" * 60)
    print(f"Output Directory: {OUTPUT_DIR}\n")

    for filename, generator, name in templates:
        try:
            prs = generator()
            filepath = os.path.join(OUTPUT_DIR, filename)
            prs.save(filepath)
            print(f"✓ Created: {name} ({filename})")
        except Exception as e:
            print(f"✗ Failed: {name} - {str(e)}")

    print("\n" + "=" * 60)
    print(f"Generated {len(templates)} templates successfully!")
    print("=" * 60)


if __name__ == "__main__":
    main()
